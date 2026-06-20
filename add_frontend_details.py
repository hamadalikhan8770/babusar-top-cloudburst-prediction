"""
Add detailed frontend customization and HTML/CSS/JavaScript implementation details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(25, 45, 85)
LIGHT_BLUE = RGBColor(100, 150, 220)
ACCENT_ORANGE = RGBColor(255, 140, 0)
GREEN = RGBColor(50, 180, 100)
RED = RGBColor(220, 80, 80)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 80, 80)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_background_gradient(slide):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = RGBColor(200, 220, 240)
    fill.gradient_stops[1].color.rgb = RGBColor(240, 248, 255)

def add_title_bar(slide, title):
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = ACCENT_ORANGE
    title_shape.line.width = Pt(2)
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(5)
    p.space_after = Pt(5)

prs = Presentation(r'C:\Users\HYPER Computers\Downloads\a4.project.bilstm\Cloud_Burst_Prediction_Presentation.pptx')

# Slide: HTML Structure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Frontend: HTML Structure (index.html)")

content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = content.text_frame
tf.word_wrap = True

html_content = [
    "KEY HTML COMPONENTS IN templates/index.html:",
    "",
    "1. PREDICTION DISPLAY SECTION",
    "   <div id='prediction-box'>",
    "     <h2>Cloudburst Prediction</h2>",
    "     <div id='probability'>72%</div>",
    "     <div id='risk-level'>HIGH RISK</div>",
    "   </div>",
    "",
    "2. REAL-TIME METRICS DISPLAY",
    "   <div id='metrics-container'>",
    "     <div class='metric-card'> CAPE: 2450 J/kg</div>",
    "     <div class='metric-card'> RH@700hPa: 78%</div>",
    "     <div class='metric-card'> Vertical Velo: 0.15 Pa/s</div>",
    "   </div>",
    "",
    "3. PREDICTION HISTORY CHART",
    "   <canvas id='chart' width='400' height='200'></canvas>",
    "",
    "4. FOOTER & SYSTEM INFO",
    "   <footer id='footer'>",
    "     Last updated: {{timestamp}}",
    "     Model Status: ✓ Ready",
    "   </footer>",
]

for idx, item in enumerate(html_content):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(9.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(0.5)
    p.space_after = Pt(0.5)
    if "<" in item and ">" in item:
        p.font.name = 'Courier New'

# Slide: CSS Styling & Cloud Theme
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "CSS Styling: Cloud-Themed Design")

left_css = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.7), Inches(6))
tf = left_css.text_frame
tf.word_wrap = True

css_design = [
    "CUSTOM CSS STYLES:",
    "",
    "1. COLOR SCHEME",
    "   • Primary: #1A2D55 (Dark storm blue)",
    "   • Accent: #FF8C00 (Thunder orange)",
    "   • Background gradient:",
    "     linear-gradient(135deg,",
    "     #C8DCF0 0%, #F0F8FF 100%)",
    "",
    "2. BACKGROUND ELEMENTS",
    "   • SVG cloud shapes",
    "   • Animated rain effect",
    "   • Thunder bolt decorations",
    "",
    "3. TYPOGRAPHY",
    "   • Font: Poppins, sans-serif",
    "   • Heading: Bold 32px",
    "   • Body: Regular 14px",
    "",
    "4. RESPONSIVE DESIGN",
    "   • Mobile: flex layout",
    "   • Tablet: 2-column grid",
    "   • Desktop: 3-column layout",
]

for idx, item in enumerate(css_design):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)

right_css = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.3), Inches(6))
tf = right_css.text_frame
tf.word_wrap = True

css_effects = [
    "VISUAL EFFECTS:",
    "",
    "5. PREDICTION BOX ANIMATION",
    "   • Pulse effect when",
    "     probability > 0.65",
    "   • Color changes:",
    "     Green: <30% (Low)",
    "     Yellow: 30-65% (Med)",
    "     Red: >65% (High)",
    "",
    "6. SHADOW & DEPTH",
    "   • Box-shadow effects",
    "   • Elevated metric cards",
    "   • Border-radius: 12px",
    "",
    "7. TRANSITIONS",
    "   • Smooth 0.3s color changes",
    "   • Chart animations",
    "   • Hover effects on buttons",
]

for idx, item in enumerate(css_effects):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_BLUE
    p.space_before = Pt(1)
    p.space_after = Pt(1)

# Slide: JavaScript Implementation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "JavaScript: Real-Time Updates & Interactivity")

js_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = js_box.text_frame
tf.word_wrap = True

js_code = [
    "KEY JAVASCRIPT FUNCTIONS:",
    "",
    "1. FETCH LIVE PREDICTION (Auto-refresh every hour)",
    "   function updatePrediction() {",
    "     fetch('/api/predict_live', {method: 'POST', headers: {'X-API-Key': API_KEY}})",
    "     .then(r => r.json())",
    "     .then(data => {",
    "       document.getElementById('probability').textContent = (data.probability*100).toFixed(1)+'%';",
    "       updateRiskColor(data.probability);",
    "       addToChart(data);",
    "     });",
    "   }",
    "",
    "2. UPDATE VISUAL INDICATORS",
    "   function updateRiskColor(prob) {",
    "     const riskBox = document.getElementById('risk-level');",
    "     if (prob > 0.65) riskBox.style.backgroundColor = '#DC5050';  // Red",
    "     else if (prob > 0.30) riskBox.style.backgroundColor = '#FFD700';  // Yellow",
    "     else riskBox.style.backgroundColor = '#32B464';  // Green",
    "   }",
    "",
    "3. RENDER PREDICTION HISTORY CHART (Chart.js)",
    "   const ctx = document.getElementById('chart').getContext('2d');",
    "   const chart = new Chart(ctx, {type: 'line', data: predictionHistory, ...});",
    "",
    "4. AUTO-REFRESH INTERVAL",
    "   setInterval(updatePrediction, 3600000);  // Every 1 hour",
]

for idx, item in enumerate(js_code):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(8.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(0.5)
    p.space_after = Pt(0.5)
    if "{" in item or "}" in item or "function" in item:
        p.font.name = 'Courier New'

# Slide: Integration Points
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Integration Points: Frontend ↔ Backend")

left_integ = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.7), Inches(6))
tf = left_integ.text_frame
tf.word_wrap = True

integ_points = [
    "FRONTEND → BACKEND:",
    "",
    "1. Request Prediction",
    "   POST /api/predict_live",
    "   Headers: X-API-Key",
    "   Body: {era5: {rows: [...],",
    "          timesteps: 6}}",
    "",
    "2. Request History",
    "   GET /api/history?days=30",
    "",
    "3. Request Model Status",
    "   GET /api/status",
    "   Returns: {",
    "     model_loaded: true,",
    "     scaler_loaded: true,",
    "     last_update: timestamp",
    "   }",
]

for idx, item in enumerate(integ_points):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.space_before = Pt(0.5)
    p.space_after = Pt(0.5)

right_integ = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.3), Inches(6))
tf = right_integ.text_frame
tf.word_wrap = True

response_data = [
    "BACKEND → FRONTEND:",
    "",
    "Response Format:",
    "{",
    "  'probability': 0.72,",
    "  'risk_level': 'HIGH',",
    "  'timestamp': '2024-...',",
    "  'metrics': {",
    "    'cape': 2450,",
    "    'rh_700': 78.5,",
    "    'vv': 0.15,",
    "    't2m': 290.2",
    "  },",
    "  'confidence': 0.94,",
    "  'lead_hours': 2",
    "}",
]

for idx, item in enumerate(response_data):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_BLUE
    p.space_before = Pt(0.5)
    p.space_after = Pt(0.5)

# Slide: Modifications & Customization
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "How to Modify & Customize the Dashboard")

modify_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = modify_box.text_frame
tf.word_wrap = True

modifications = [
    "COMMON CUSTOMIZATIONS:",
    "",
    "1. CHANGE ALERT THRESHOLDS",
    "   Edit app.py: PROBABILITY_THRESHOLD = 0.65 → Change to 0.75 for stricter alerts",
    "",
    "2. UPDATE DASHBOARD COLORS",
    "   Edit templates/index.html CSS:",
    "   :root { --primary: #1A2D55; --accent: #FF8C00; }",
    "   Change hex values to desired colors",
    "",
    "3. ADD NEW METRICS DISPLAY",
    "   Add to index.html: <div class='metric-card' id='new-metric'></div>",
    "   Update JavaScript: document.getElementById('new-metric').textContent = data.new_metric;",
    "",
    "4. MODIFY REFRESH RATE",
    "   app.py: setInterval(updatePrediction, 3600000) → Change milliseconds",
    "   3600000ms = 1 hour, 600000ms = 10 minutes, 60000ms = 1 minute",
    "",
    "5. CHANGE LOCATION (Babusar Top → Another Location)",
    "   app.py: LATITUDE = 35.6, LONGITUDE = 73.6 → Update coordinates",
    "   Database stores location-specific models and scalers",
    "",
    "6. ADD EMAIL/SMS ALERTS",
    "   app.py: When probability > threshold, call send_email() or send_sms()",
    "   Install: pip install twilio, pip install Flask-Mail",
]

for idx, item in enumerate(modifications):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.space_before = Pt(0.8)
    p.space_after = Pt(0.8)
    if "Edit" in item or "Add" in item or "Change" in item:
        p.font.bold = True

# Slide: Deployment Instructions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Deployment: From Development to Production")

deploy_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = deploy_box.text_frame
tf.word_wrap = True

deployment = [
    "DEVELOPMENT ENVIRONMENT:",
    "  1. pip install -r requirements.txt",
    "  2. python app.py",
    "  3. Open http://localhost:5000 in browser",
    "",
    "PRODUCTION DEPLOYMENT OPTIONS:",
    "",
    "Option 1: LOCAL SERVER (Windows/Linux)",
    "  • Install Gunicorn: pip install gunicorn",
    "  • Run: gunicorn -w 4 -b 0.0.0.0:5000 app:app",
    "  • Configure reverse proxy (Nginx)",
    "",
    "Option 2: DOCKER CONTAINERIZATION",
    "  • Create Dockerfile with all dependencies",
    "  • Build: docker build -t cloudburst-predictor .",
    "  • Run: docker run -p 5000:5000 cloudburst-predictor",
    "",
    "Option 3: CLOUD DEPLOYMENT",
    "  • AWS Lambda/EC2: Upload app.py + requirements.txt + model files",
    "  • Google Cloud Run: Deploy from git repository",
    "  • Heroku: git push heroku main (auto-deployment)",
    "",
    "Option 4: API SERVER (Production Grade)",
    "  • Use uWSGI: pip install uwsgi",
    "  • Nginx as reverse proxy for SSL/TLS encryption",
    "  • Enable HTTPS for security (free SSL: Let's Encrypt)",
]

for idx, item in enumerate(deployment):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(9.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(0.5)
    p.space_after = Pt(0.5)
    if "Option" in item or "ENVIRONMENT" in item or "OPTIONS" in item:
        p.font.bold = True

# Slide: Testing & Performance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Testing, Monitoring & Maintenance")

test_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = test_box.text_frame
tf.word_wrap = True

testing = [
    "TESTING THE DASHBOARD:",
    "",
    "1. UNIT TESTS (Test individual functions)",
    "   • Test feature extraction: Does it correctly build (6, 26) tensor?",
    "   • Test model inference: Does BiLSTM return probability in [0, 1]?",
    "   • Test API endpoints: Does /api/predict_live return valid JSON?",
    "",
    "2. INTEGRATION TESTS",
    "   • Full pipeline: Open-Meteo → feature extraction → model → dashboard",
    "   • Check data freshness: Is latest prediction within 1 hour?",
    "",
    "3. PERFORMANCE MONITORING",
    "   • Response time: Target <500ms per request (measure with time.time())",
    "   • Memory usage: BiLSTM inference should use <200MB RAM",
    "   • CPU load: Monitor during peak prediction requests",
    "",
    "4. PRODUCTION MONITORING",
    "   • Error logging: Log all API failures with timestamp",
    "   • Prediction accuracy: Track actual vs predicted outcomes",
    "   • Uptime monitoring: Use tools like Uptime Robot",
    "",
    "5. MAINTENANCE TASKS",
    "   • Weekly: Check model performance, false alarm rate",
    "   • Monthly: Update dependencies (pip list --outdated)",
    "   • Quarterly: Retrain model with new data if available",
]

for idx, item in enumerate(testing):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(9.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(0.5)
    p.space_after = Pt(0.5)
    if ":" in item and item[0].isdigit():
        p.font.bold = True

prs.save(r'C:\Users\HYPER Computers\Downloads\a4.project.bilstm\Cloud_Burst_Prediction_Presentation.pptx')
print("✓ Frontend customization slides added!")
print(f"✓ Total slides: {len(prs.slides)}")
print("✓ Added: HTML structure, CSS styling, JavaScript implementation")
print("✓ Added: Integration points, Modifications guide, Deployment options")
print("✓ Added: Testing and maintenance guidelines")
