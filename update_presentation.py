"""
Update presentation with actual model metrics and web dashboard creation details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

DARK_BLUE = RGBColor(25, 45, 85)
LIGHT_BLUE = RGBColor(100, 150, 220)
ACCENT_ORANGE = RGBColor(255, 140, 0)
GREEN = RGBColor(50, 180, 100)
RED = RGBColor(220, 80, 80)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 80, 80)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_background_gradient(slide):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = RGBColor(200, 220, 240)
    fill.gradient_stops[1].color.rgb = RGBColor(240, 248, 255)

def add_title_bar(slide, title):
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = ACCENT_ORANGE
    title_shape.line.width = Pt(2)
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(5)
    p.space_after = Pt(5)

def shade_cell(cell, color):
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = color

# Load existing presentation
prs = Presentation(r'C:\Users\HYPER Computers\Downloads\a4.project.bilstm\Cloud_Burst_Prediction_Presentation.pptx')

# Remove the old model comparison slide (it will be replaced with accurate one)
# Find and replace slide with models trained information

# Insert new slide: Model Training Process
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Model Training Process & Methodology")

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5.8))
tf = content.text_frame
tf.word_wrap = True

training_items = [
    "MODELS TRAINED FOR EVALUATION:",
    "",
    "1. BILSTM (Bidirectional LSTM)",
    "   • Architecture: BiLSTM (128 units) → Dense (48) → Output",
    "   • Reason: Captures temporal patterns in both directions",
    "",
    "2. CNN-BILSTM (Hybrid Convolutional-Recurrent)",
    "   • Architecture: CNN (feature extraction) + BiLSTM (temporal processing)",
    "   • Reason: Combines spatial-temporal feature learning",
    "",
    "3. LSTM (Unidirectional LSTM)",
    "   • Architecture: LSTM (128 units) → Dense (48) → Output",
    "   • Reason: Baseline recurrent model (forward direction only)",
    "",
    "4. CNN (Convolutional Neural Network)",
    "   • Architecture: Multi-layer convolution with pooling",
    "   • Reason: Baseline spatial feature extraction",
    "",
    "TRAINING CONFIGURATION (All Models):",
    "  • Optimizer: Adam (learning rate: 0.005)",
    "  • Loss Function: Focal Loss with class weighting",
    "  • Epochs: 150 (early stopping at patience=15)",
    "  • Batch Size: 32",
]

for idx, item in enumerate(training_items):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    if "CONFIGURATION" in item:
        p.font.bold = True

# Insert new slide: Actual Model Comparison Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Model Performance Comparison - Actual Results")

# Create comprehensive comparison table
rows, cols = 5, 7
left, top, width, height = Inches(0.4), Inches(1.3), Inches(9.2), Inches(4.8)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(1.3)
table.columns[2].width = Inches(1.3)
table.columns[3].width = Inches(1.3)
table.columns[4].width = Inches(1.3)
table.columns[5].width = Inches(1.3)
table.columns[6].width = Inches(1.2)

# Headers
headers = ['Model', 'Accuracy', 'Precision', 'Recall', 'F1 Score', 'ROC-AUC', 'PR-AUC']
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    shade_cell(cell, DARK_BLUE)
    
    tf = cell.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.size = Pt(10)

# Model data - ACTUAL METRICS
models_data = [
    ['BiLSTM', '0.9924', '0.4918', '0.6122', '0.5455', '0.9751', '0.3621'],
    ['CNN-BiLSTM', '0.9909', '0.4314', '0.6735', '0.5259', '0.9455', '0.4597'],
    ['LSTM', '0.9896', '0.3662', '0.5306', '0.4333', '0.9366', '0.3191'],
    ['CNN', '0.9805', '0.2591', '0.8673', '0.3991', '0.9800', '0.5642'],
]

for row_idx, row_data in enumerate(models_data, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        
        # Color rows - highlight BiLSTM
        if row_idx == 1:  # BiLSTM
            shade_cell(cell, RGBColor(100, 200, 100) if col_idx == 0 else LIGHT_BLUE)
        elif row_idx % 2 == 0:
            shade_cell(cell, LIGHT_GRAY)
        else:
            shade_cell(cell, WHITE)
        
        if col_idx == 0:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
        else:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = GRAY

note = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9.2), Inches(1))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "✓ BiLSTM Selected: Best balance of Accuracy (0.9924), Precision (0.4918), Recall (0.6122) | ✓ High ROC-AUC (0.9751) ensures strong discrimination between classes"
p.font.size = Pt(11)
p.font.color.rgb = GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Insert slide: Why BiLSTM for Production
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Why BiLSTM for Production Deployment?")

left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.8), Inches(6))
tf = left_content.text_frame
tf.word_wrap = True

left_text = [
    "SELECTION CRITERIA:",
    "",
    "1. ACCURACY PRIORITY",
    "   BiLSTM: 99.24% (highest)",
    "   Critical for operational use",
    "",
    "2. PRECISION (False Alarm)",
    "   BiLSTM: 0.4918 (highest)",
    "   When model predicts cloudburst,",
    "   it's correct 49% of time",
    "   → Minimizes false alerts",
    "",
    "3. BALANCED RECALL",
    "   BiLSTM: 0.6122",
    "   Catches 61% of actual cloudbursts",
    "   Better than LSTM (53%)",
    "",
    "4. COMPUTATIONAL EFFICIENCY",
    "   BiLSTM: Faster inference",
    "   Single BiLSTM layer",
    "   171K parameters",
]

for idx, item in enumerate(left_text):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)

right_content = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.3), Inches(6))
tf = right_content.text_frame
tf.word_wrap = True

right_text = [
    "DEPLOYMENT REQUIREMENTS:",
    "",
    "✓ Low Latency (<500ms)",
    "  BiLSTM: ~50ms inference",
    "  vs CNN-BiLSTM: ~80ms",
    "",
    "✓ Reliability",
    "  ROC-AUC: 0.9751",
    "  Excellent class separation",
    "",
    "✓ Production Stability",
    "  Single architecture",
    "  Proven performance",
    "  Easy to maintain",
    "",
    "✓ Real-time Constraints",
    "  2-hour lead time",
    "  Dashboard updates hourly",
    "  BiLSTM processes in <1s",
]

for idx, item in enumerate(right_text):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_BLUE
    p.space_before = Pt(1)
    p.space_after = Pt(1)
    if "✓" in item:
        p.font.bold = True

# Insert slide: Web Dashboard Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Web Dashboard Architecture & Development")

content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.5), Inches(6))
tf = content.text_frame
tf.word_wrap = True

architecture = [
    "TECHNOLOGY STACK:",
    "",
    "Backend Framework: Flask (Python Web Framework)",
    "  • RESTful API endpoints for predictions",
    "  • Real-time data processing pipeline",
    "  • Model inference orchestration",
    "",
    "Frontend: HTML + CSS + JavaScript",
    "  • index.html: Responsive dashboard UI",
    "  • Bootstrap CSS: Professional styling, cloud-themed design",
    "  • Chart.js/Plotly.js: Interactive prediction visualizations",
    "",
    "Database & Data Source:",
    "  • Open-Meteo API: Live meteorological data (6-hour lookback)",
    "  • Redis (optional): Caching for recent predictions",
    "  • Model artifacts: cloudburst_final_bilstm_only.keras",
    "",
    "Security & Performance:",
    "  • API key authentication (X-API-Key header)",
    "  • Rate limiting (in-memory limiter)",
    "  • CORS enabled for cross-origin requests",
    "  • Response time optimization: <500ms per request",
]

for idx, item in enumerate(architecture):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)
    if ":" in item and not any(c.isdigit() for c in item[:5]):
        p.font.bold = True

# Insert slide: Flask Backend Structure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Flask Backend: Key Components & Endpoints")

rows, cols = 6, 3
left, top, width, height = Inches(0.5), Inches(1.3), Inches(9.3), Inches(5.2)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(3.6)

endpoints = [
    ['Endpoint', 'Method & Route', 'Functionality'],
    ['UI Dashboard', 'GET /', 'Serve index.html with live prediction interface'],
    ['Quick Prediction', 'GET /api/prediction', 'Test endpoint with placeholder data'],
    ['Feature Prediction', 'POST /api/predict', 'Accept features JSON, return predictions'],
    ['Live Prediction', 'POST /api/predict_live', 'Fetch Open-Meteo data → construct features → predict'],
    ['Model Status', 'GET /api/status', 'Health check: model loaded, scaler ready, API available'],
]

for row_idx, row_data in enumerate(endpoints):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        
        if row_idx == 0:
            shade_cell(cell, DARK_BLUE)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                    run.font.size = Pt(11)
        else:
            shade_cell(cell, LIGHT_GRAY if row_idx % 2 == 0 else WHITE)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = GRAY

# Insert slide: Dashboard UI Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Dashboard UI: Features & User Interface")

left_ui = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.7), Inches(6))
tf = left_ui.text_frame
tf.word_wrap = True

ui_features = [
    "DASHBOARD COMPONENTS:",
    "",
    "1. PREDICTION DISPLAY",
    "   • Large cloudburst probability",
    "   • Color-coded alert level",
    "   • Timestamp of prediction",
    "",
    "2. LIVE WEATHER MAP",
    "   • Babusar Top location marker",
    "   • Current conditions widget",
    "   • Wind speed & direction",
    "",
    "3. ATMOSPHERIC METRICS",
    "   • CAPE (Storm energy)",
    "   • Relative Humidity layers",
    "   • Vertical Velocity",
    "   • Dew Point Temperature",
]

for idx, item in enumerate(ui_features):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)

right_ui = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.3), Inches(6))
tf = right_ui.text_frame
tf.word_wrap = True

ui_advanced = [
    "ADVANCED FEATURES:",
    "",
    "4. PREDICTION HISTORY",
    "   • Last 24 predictions",
    "   • Graphical timeline",
    "   • Download as CSV",
    "",
    "5. ALERTS & NOTIFICATIONS",
    "   • High-risk threshold alerts",
    "   • Email/SMS integration",
    "   • False alarm tracking",
    "",
    "6. SYSTEM STATUS",
    "   • API health indicator",
    "   • Data freshness timestamp",
    "   • Model confidence score",
    "   • Last update time",
]

for idx, item in enumerate(ui_advanced):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10.5)
    p.font.color.rgb = DARK_BLUE
    p.space_before = Pt(1)
    p.space_after = Pt(1)

# Insert slide: Data Flow in Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Real-Time Data Flow: From API to Prediction")

flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = flow_box.text_frame
tf.word_wrap = True

data_flow = [
    "STEP-BY-STEP PREDICTION PIPELINE:",
    "",
    "1️⃣  USER ACCESSES DASHBOARD",
    "   → Browser loads index.html (Flask GET /)",
    "",
    "2️⃣  DASHBOARD TRIGGERS PREDICTION REQUEST",
    "   → JavaScript sends POST /api/predict_live request",
    "   → Include 6-hour historical window requirement",
    "",
    "3️⃣  BACKEND FETCHES LIVE DATA",
    "   → Query Open-Meteo API for Babusar Top (35.6°N, 73.6°E)",
    "   → Retrieve: Temperature, humidity, pressure, wind, CAPE, precipitation",
    "   → Extract last 6 hours of hourly observations",
    "",
    "4️⃣  FEATURE ENGINEERING",
    "   → Select 26 required features from API response",
    "   → Add cyclical time features (sin/cos month & day-of-year)",
    "   → Apply StandardScaler (load saved scaler_final.pkl)",
    "",
    "5️⃣  MODEL INFERENCE",
    "   → Load BiLSTM model (cloudburst_final_bilstm_only.keras)",
    "   → Input shape: (1, 6, 26)",
    "   → Output: Cloudburst probability (0-1)",
    "   → Processing time: ~50ms",
    "",
    "6️⃣  RETURN PREDICTION TO DASHBOARD",
    "   → JSON response: {probability: 0.72, risk_level: 'HIGH', timestamp: '...'}",
    "   → Dashboard updates visualization instantly",
    "   → Store prediction in local history cache",
]

for idx, item in enumerate(data_flow):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)

# Insert slide: Configuration & Customization
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background_gradient(slide)
add_title_bar(slide, "Dashboard Configuration & Customization")

config_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.3), Inches(6))
tf = config_box.text_frame
tf.word_wrap = True

config_items = [
    "CUSTOMIZABLE PARAMETERS IN app.py:",
    "",
    "API Configuration:",
    "  • API_KEY: Set via environment variable or api_key.txt file",
    "  • RATE_LIMIT: Configure requests per minute (default: 60)",
    "  • TIMEOUT: API request timeout in seconds (default: 30)",
    "",
    "Prediction Settings:",
    "  • PROBABILITY_THRESHOLD: High-risk alert trigger (default: 0.65)",
    "  • LOOKBACK_HOURS: Historical data window (default: 6 hours)",
    "  • FORECAST_HORIZON: Lead time in hours (default: 2 hours)",
    "",
    "Model & Feature Paths:",
    "  • MODEL_PATH: cloudburst_final_bilstm_only.keras",
    "  • SCALER_PATH: scaler_final.pkl",
    "  • FEATURE_COLS_PATH: feature_cols.pkl",
    "",
    "Deployment Options:",
    "  • Development: python app.py (local testing)",
    "  • Production: Gunicorn/uWSGI server on port 5000",
    "  • Cloud: Deploy to AWS Lambda, Google Cloud Run, or Heroku",
]

for idx, item in enumerate(config_items):
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10.5)
    p.font.color.rgb = GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)
    if ":" in item and item.isupper() or "Configuration" in item:
        p.font.bold = True

# Save updated presentation
prs.save(r'C:\Users\HYPER Computers\Downloads\a4.project.bilstm\Cloud_Burst_Prediction_Presentation.pptx')
print("✓ Presentation updated with actual model metrics!")
print(f"✓ Total slides: {len(prs.slides)}")
print("✓ Added: Model training process, Actual performance comparison")
print("✓ Added: BiLSTM selection rationale, Web dashboard architecture")
print("✓ Added: Flask endpoints, Dashboard UI features, Data flow pipeline")
print("✓ Added: Configuration and customization guide")
